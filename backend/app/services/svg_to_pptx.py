"""
SVG 转 PPTX 导出器
将 AI 生成的 SVG 页面转换为原生 PowerPoint 文件

方案：使用 svglib 将 SVG 转为高清 PNG，再嵌入 PPTX
优势：
- 高清图片，显示清晰
- 兼容所有版本的 PowerPoint
- 文件大小适中
"""

import os
import tempfile
from pathlib import Path
from typing import List, Optional


def convert_svgs_to_pptx(svg_dir: str, output_path: str, design_spec: str = ""):
    """
    将目录下的所有 SVG 文件转换为 PPTX

    Args:
        svg_dir: SVG 文件所在目录
        output_path: 输出的 PPTX 文件路径
        design_spec: 设计规范内容（可选，用于提取标题）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu

        prs = Presentation()

        # 设置 16:9 宽屏尺寸（13.333 x 7.5 英寸）
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 获取所有 SVG 文件并按页码排序
        svg_files = sorted(Path(svg_dir).glob("*.svg"),
                          key=lambda f: _extract_page_number(f.name))

        if not svg_files:
            raise ValueError(f"在 {svg_dir} 中未找到 SVG 文件")

        # 提取 PPT 标题
        ppt_title = _extract_title_from_spec(design_spec) or "未命名演示文稿"

        # 添加封面页（从第一个 SVG 提取）
        if svg_files:
            _add_cover_slide(prs, svg_files[0], ppt_title)

        # 添加内容页
        for svg_file in svg_files:
            _add_svg_slide(prs, svg_file)

        # 保存 PPTX
        prs.save(output_path)
        print(f"PPTX 导出成功: {output_path} ({len(svg_files)} 页)")

    except ImportError as e:
        raise ImportError(f"缺少必要依赖: {e}。请运行: pip install python-pptx")

    except Exception as e:
        raise RuntimeError(f"PPTX 导出失败: {e}")


def _svg_to_png(svg_path: str, dpi: int = 150) -> str:
    """
    将 SVG 转为 PNG（使用 PyMuPDF）

    Args:
        svg_path: SVG 文件路径
        dpi: 输出分辨率

    Returns:
        临时 PNG 文件路径
    """
    import fitz  # PyMuPDF

    # 打开 SVG 文件
    doc = fitz.open(svg_path)

    # 计算缩放比例（96 DPI 是 SVG 默认，转为目标 DPI）
    zoom = dpi / 96
    mat = fitz.Matrix(zoom, zoom)

    # 渲染第一页为图片
    page = doc[0]
    pix = page.get_pixmap(matrix=mat)

    # 创建临时 PNG 文件
    fd, png_path = tempfile.mkstemp(suffix='.png')
    os.close(fd)
    pix.save(png_path)

    doc.close()

    return png_path


def _add_svg_slide(prs, svg_file: Path):
    """
    添加一页 SVG 内容
    先将 SVG 转为高清 PNG，再嵌入 PPTX
    """
    from pptx.util import Inches

    # 创建空白幻灯片
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)

    # 将 SVG 转为 PNG
    png_path = _svg_to_png(str(svg_file), dpi=150)

    try:
        # 插入 PNG 图片填满幻灯片
        slide.shapes.add_picture(
            png_path,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
    finally:
        # 清理临时 PNG
        try:
            os.remove(png_path)
        except:
            pass


def _add_cover_slide(prs, svg_file: Path, title: str):
    """
    添加封面页
    使用第一张 SVG 作为封面背景，并添加标题
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 将 SVG 转为 PNG 并添加为背景
    png_path = _svg_to_png(str(svg_file), dpi=150)

    try:
        slide.shapes.add_picture(
            png_path,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
    finally:
        try:
            os.remove(png_path)
        except:
            pass

    # 在封面上添加标题文本框
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(2)
    )

    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Microsoft YaHei"

    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(12.333), Inches(1)
    )

    tf2 = subtitle_box.text_frame
    tf2.text = "AI 智能生成"
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.font.name = "Microsoft YaHei"


def _extract_page_number(filename: str) -> int:
    """从文件名提取页码数字"""
    import re
    match = re.search(r'(\d+)', filename)
    return int(match.group(1)) if match else 0


def _extract_title_from_spec(design_spec: str) -> Optional[str]:
    """从设计规范中提取标题"""
    if not design_spec:
        return None

    import re
    # 匹配 YAML frontmatter 中的 title
    match = re.search(r'title:\s*"([^"]+)"', design_spec)
    if match:
        return match.group(1)

    # 匹配 Markdown 一级标题
    match = re.search(r'^#\s+(.+)$', design_spec, re.MULTILINE)
    if match:
        return match.group(1)

    return None


def get_pptx_info(pptx_path: str) -> dict:
    """
    获取 PPTX 文件信息

    Returns:
        {"pages": 页数, "title": 标题, "size_kb": 文件大小}
    """
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        file_size = os.path.getsize(pptx_path) / 1024

        # 尝试提取标题（从第一个文本框）
        title = None
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and shape.text:
                    title = shape.text.strip()
                    break

        return {
            "pages": len(prs.slides),
            "title": title or "未命名",
            "size_kb": round(file_size, 1)
        }
    except Exception as e:
        return {"pages": 0, "title": "未知", "size_kb": 0, "error": str(e)}
